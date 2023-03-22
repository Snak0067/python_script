# -*- coding:utf-8 -*-
# @FileName :api.py
# @Time :2023/3/21 17:00
# @Author :Xiaofeng
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Cm


def loadcsv():
    """
    读取正确词目文件及错误词目csv文件，并转成字典
    :return: id_to_gloss, gloss_to_id
    """
    word_reader = pd.read_csv('rightForm.csv', encoding='utf-8')
    words = np.array(word_reader)
    return words


def generate(words: list):
    """
    生成ppt文本
    :param words: ppt内容的字典数据
    :return:
    """
    # 创建空白演示文稿
    prs = Presentation()
    for data in words:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        # 在指定位置添加文本框
        textbox = slide.shapes.add_textbox(Cm(3), Cm(5), Cm(10), Cm(10))
        tf = textbox.text_frame
        # 在文本框中写入文字
        para_title = tf.add_paragraph()  # 新增段落
        para_title.text = data[0]  # 向段落写入文字
        # 设置字体
        font = para_title.font
        font.name = '微软雅黑'  # 字体类型
        font.bold = True  # 加粗
        font.size = Pt(64)  # 大小

        # 在指定位置添加文本框
        textbox = slide.shapes.add_textbox(Cm(3), Cm(2), Cm(10), Cm(10))
        tf = textbox.text_frame
        # 在文本框中写入文字
        para_title = tf.add_paragraph()  # 新增段落
        para_title.text = str(data[1])  # 向段落写入文字
        # 设置字体
        font = para_title.font
        font.name = '微软雅黑'  # 字体类型
        font.bold = True  # 加粗
        font.size = Pt(64)  # 大小
    # 保存
    prs.save('new_medical_dict.pptx')


if __name__ == '__main__':
    words = loadcsv()
    generate(words)